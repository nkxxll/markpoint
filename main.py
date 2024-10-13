from argparse import ArgumentParser
from os.path import isfile
from sys import stdin

from pptx import Presentation

from markpoint import MarkParser, PointGenerator


def main():
    parser = ArgumentParser()
    parser.add_argument("markdown", type=str, default="", help="Markdown file")
    parser.add_argument(
        "-o", "--outfile", type=str, help="Outfile for the powerpoint presentation"
    )
    parser.add_argument(
        "-i",
        "--inputfile",
        type=str,
        default=None,
        help="Input markdown file for the powerpoint presentation",
    )
    parser.add_argument(
        "-a",
        "--at",
        type=int,
        required=False,
        help="Insert markdown generate slides at zero based idx",
    )
    args = parser.parse_args()

    if args.markdown == "":
        text = stdin.read()
    else:
        if not isfile(args.markdown):
            raise Exception(f"Markdown file {args.markdown} could not be found")
        with open(args.markdown, "r") as f:
            text = f.read()

    prs = None
    if args.inputfile is not None:
        prs = Presentation(args.inputfile)
    else:
        prs = Presentation()

    mp = MarkParser(text)
    pg = PointGenerator(mp.slides, prs)
    # todo insert at
    pg.point_from_slideinfos()
    pg.save(args.outfile)


if __name__ == "__main__":
    main()
