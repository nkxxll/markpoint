from enum import Enum, auto

from pptx import Presentation

from ._slide import SlideInfo

__all__ = ["PointGenerator"]


class Layouts(Enum):
    TITLE = 0
    TITLE_CONTENT = auto
    SECTION_HEADER = auto
    TWO_CONTENT = auto
    COMPARISON = auto
    TITLE_ONLY = auto
    BLANK = auto
    CONTENT_CAPTION = auto
    PICTURE_CAPTION = auto


class PointGenerator:
    def __init__(self, slideinfo: list[SlideInfo], prs) -> None:
        self._slideinfolist = slideinfo
        self._prs = prs

    @property
    def slideinfo_list(self):
        return self._slideinfolist

    def insert_slide(self, slideInfo: SlideInfo, index, layout=Layouts):
        # Add the new slide
        new_slide = self._prs.slides.add_slide(self._prs.slide_layouts[layout])
        new_slide.shapes.title.text = slideInfo.header

        # Temporarily store all slides in the desired order
        slides = list(self._prs.slides._sldIdLst)
        self._prs.slides._sldIdLst.remove(slides[-1])  # Remove the new slide
        slides.insert(index, slides.pop(-1))  # Insert it at the specified index

        # Apply the new order to the presentation
        self._prs.slides._sldIdLst[:] = slides

    def point_from_slideinfos(self) -> None:
        """generate powerpoint presentation

        Args:
            outfile: file to save the powerpoint presentation at
        """
        # Create a presentation object
        for s in self._slideinfolist:
            # Add a blank slide layout
            slide_layout = self._prs.slide_layouts[1]  # Title and Content layout
            slide = self._prs.slides.add_slide(slide_layout)

            # Set the title
            title = slide.shapes.title
            title.text = s.header

            # Add bullet points to the content placeholder
            content = slide.placeholders[1]
            for t in s.texts:
                p = content.text_frame.add_paragraph()
                p.text = t.text
                p.level = t.level

    def save(self, outfile: str) -> None:
        self._prs.save(outfile)
