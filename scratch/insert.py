from pptx import Presentation
from pptx.oxml import slide


# Function to copy a slide from one presentation to another
def copy_slide(prs, slide):
    new_slide = prs.slides.add_slide(slide.slide_layout)

    # Copy all shapes from the original slide
    for shape in slide.shapes:
        if shape.has_text_frame:
            textbox = new_slide.shapes.add_textbox(
                shape.left, shape.top, shape.width, shape.height
            )
            textbox.text = shape.text

    return new_slide


# Create original presentation
original_prs = Presentation("./test.pptx")
slides = original_prs.slides


# Save the new presentation
new_prs.save("ordered_presentation.pptx")
