from pptx import Presentation

# Initialize the presentation
prs = Presentation("test.pptx")


# Function to insert a new slide at a specific index
def insert_slide(prs, index, layout=0, title="Inserted Slide"):
    # Add the new slide
    new_slide = prs.slides.add_slide(prs.slide_layouts[layout])
    new_slide.shapes.title.text = title

    # Temporarily store all slides in the desired order
    slides = list(prs.slides._sldIdLst)
    prs.slides._sldIdLst.remove(slides[-1])  # Remove the new slide
    slides.insert(index, slides.pop(-1))  # Insert it at the specified index

    # Apply the new order to the presentation
    prs.slides._sldIdLst[:] = slides


# Insert a slide at index 1 (second position)
insert_slide(prs, 0, title="Inserted Slide at Position 2")

# Save the presentation to file
prs.save("reordered_presentation.pptx")
